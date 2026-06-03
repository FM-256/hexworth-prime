#!/usr/bin/env python3
# search-all.py — Unified search across all Hexworth knowledge surfaces.
#
# Surfaces:
#   1. Repo:        /home/eq/ai-content/hexworth-prime/
#   2. Memory:      /home/eq/.claude/projects/-home-eq-ai-content-hexworth-prime/memory/
#   3. Shared:      /home/eq/hexworth-shared/   (docx/pdf/pptx/xlsx extracted, cached)
#   4. Confluence:  hexworth.atlassian.net via REST API
#
# Operational backbone for the "Search Before Asking" CLAUDE.md rule
# (added 2026-06-03). One command -> hits across every surface.
#
# Usage:
#   python3 search-all.py "query text"
#   python3 search-all.py "query" --surfaces repo,memory
#   python3 search-all.py "query" --regex
#   python3 search-all.py "query" --no-confluence
#   python3 search-all.py "query" --limit 30
#   python3 search-all.py "query" --context 2
#   python3 search-all.py "query" --build-cache

# ---------------------------------------------------------------------------
# Imports (stdlib only + shutil for which-style checks)
# ---------------------------------------------------------------------------
import argparse
import base64
import hashlib
import json
import os
import re
import shutil
import subprocess
import sys
import time
import urllib.parse
import urllib.request
from pathlib import Path

# ---------------------------------------------------------------------------
# Constants and paths
# ---------------------------------------------------------------------------
REPO_ROOT = Path("/home/eq/ai-content/hexworth-prime")
MEMORY_DIR = Path("/home/eq/.claude/projects/-home-eq-ai-content-hexworth-prime/memory")
SHARED_DIR = Path("/home/eq/hexworth-shared")
CONFLUENCE_CREDS = Path("/home/eq/.config/confluence/credentials.json")
CACHE_DIR = Path("/tmp/hexworth-search-cache")

# Repo exclusions: noisy dirs + audit-dump files
REPO_EXCLUDE_DIRS = {
    "node_modules", ".git", "_archive",
    "_tools/reports",      # huge audit JSON dumps
}
REPO_EXCLUDE_FILES = {
    "_tools/nexus/findings.json",   # giant
}
REPO_EXCLUDE_SUFFIXES = {".bak", ".backup"}

# Text extensions we grep directly (no extraction)
PLAIN_TEXT_EXTS = {
    ".md", ".txt", ".html", ".htm", ".json", ".js", ".mjs", ".cjs",
    ".ts", ".tsx", ".jsx", ".csv", ".py", ".sh", ".bash", ".zsh",
    ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf", ".env",
    ".css", ".scss", ".less", ".xml", ".rb", ".go", ".rs", ".c",
    ".cpp", ".h", ".hpp", ".java", ".kt", ".swift", ".php", ".sql",
    ".tex", ".rst", ".log", ".tsv",
}

# Extensions that need extraction
EXTRACT_EXTS = {".pdf", ".docx", ".pptx", ".xlsx"}

# Binary / never-scan suffixes (skip silently)
BINARY_SUFFIXES = {
    ".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff", ".ico",
    ".mp3", ".mp4", ".wav", ".mov", ".avi", ".mkv", ".webm", ".ogg",
    ".zip", ".tar", ".gz", ".bz2", ".xz", ".7z", ".rar",
    ".woff", ".woff2", ".ttf", ".otf", ".eot",
    ".so", ".o", ".a", ".dll", ".exe", ".dylib", ".class", ".jar",
    ".pyc", ".pyo", ".whl", ".egg",
    ".sqlite", ".db", ".dat",
}

CONFLUENCE_TIMEOUT = 10  # seconds


# ---------------------------------------------------------------------------
# Extraction tool availability check (done once at startup)
# ---------------------------------------------------------------------------
def check_tools():
    """Return dict of which extraction tools are available."""
    tools = {
        "pdftotext": shutil.which("pdftotext") is not None,
        "pandoc": shutil.which("pandoc") is not None,
    }
    try:
        import docx  # noqa: F401
        tools["python-docx"] = True
    except ImportError:
        tools["python-docx"] = False
    try:
        import pptx  # noqa: F401
        tools["python-pptx"] = True
    except ImportError:
        tools["python-pptx"] = False
    try:
        import openpyxl  # noqa: F401
        tools["openpyxl"] = True
    except ImportError:
        tools["openpyxl"] = False
    return tools


# ---------------------------------------------------------------------------
# Cache helpers (extraction is expensive; cache by path + mtime)
# ---------------------------------------------------------------------------
def cache_key_for(path: Path) -> Path:
    """Return cache file path for given source file."""
    try:
        mtime = int(path.stat().st_mtime)
    except OSError:
        mtime = 0
    h = hashlib.sha1(f"{path}|{mtime}".encode("utf-8")).hexdigest()
    return CACHE_DIR / f"{h}.txt"


def cached_text(path: Path) -> str | None:
    """Return cached extracted text, or None if cache miss / stale."""
    ck = cache_key_for(path)
    if ck.exists():
        try:
            return ck.read_text(encoding="utf-8", errors="replace")
        except OSError:
            return None
    return None


def write_cache(path: Path, text: str) -> None:
    """Write extracted text to cache."""
    CACHE_DIR.mkdir(parents=True, exist_ok=True)
    ck = cache_key_for(path)
    try:
        ck.write_text(text, encoding="utf-8")
    except OSError:
        pass


# ---------------------------------------------------------------------------
# Extractors for PDF / DOCX / PPTX / XLSX
# ---------------------------------------------------------------------------
def extract_pdf(path: Path, tools: dict) -> str | None:
    if not tools.get("pdftotext"):
        return None
    try:
        r = subprocess.run(
            ["pdftotext", "-q", "-layout", str(path), "-"],
            capture_output=True, timeout=30,
        )
        if r.returncode == 0:
            return r.stdout.decode("utf-8", errors="replace")
    except (subprocess.TimeoutExpired, OSError):
        pass
    return None


def extract_docx(path: Path, tools: dict) -> str | None:
    # Prefer pandoc (faster, handles tables) -> fallback to python-docx
    if tools.get("pandoc"):
        try:
            r = subprocess.run(
                ["pandoc", "-t", "plain", str(path)],
                capture_output=True, timeout=30,
            )
            if r.returncode == 0:
                return r.stdout.decode("utf-8", errors="replace")
        except (subprocess.TimeoutExpired, OSError):
            pass
    if tools.get("python-docx"):
        try:
            import docx
            d = docx.Document(str(path))
            return "\n".join(p.text for p in d.paragraphs)
        except Exception:
            pass
    return None


def extract_pptx(path: Path, tools: dict) -> str | None:
    if not tools.get("python-pptx"):
        return None
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        out = []
        for i, slide in enumerate(prs.slides, 1):
            out.append(f"[Slide {i}]")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text:
                                out.append(run.text)
        return "\n".join(out)
    except Exception:
        return None


def extract_xlsx(path: Path, tools: dict) -> str | None:
    if not tools.get("openpyxl"):
        return None
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
        out = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            out.append(f"[Sheet: {sheet}]")
            for row in ws.iter_rows(values_only=True):
                line = "\t".join("" if c is None else str(c) for c in row)
                if line.strip():
                    out.append(line)
        return "\n".join(out)
    except Exception:
        return None


def get_extracted_text(path: Path, tools: dict, stats: dict) -> str | None:
    """Return plain text for a non-text file, using cache when possible."""
    ext = path.suffix.lower()
    if ext not in EXTRACT_EXTS:
        return None
    cached = cached_text(path)
    if cached is not None:
        stats["cache_hits"] += 1
        return cached
    if ext == ".pdf":
        text = extract_pdf(path, tools)
        handler = "pdftotext"
    elif ext == ".docx":
        text = extract_docx(path, tools)
        handler = "pandoc/python-docx"
    elif ext == ".pptx":
        text = extract_pptx(path, tools)
        handler = "python-pptx"
    elif ext == ".xlsx":
        text = extract_xlsx(path, tools)
        handler = "openpyxl"
    else:
        return None
    if text is None:
        stats["extract_failed"][handler] = stats["extract_failed"].get(handler, 0) + 1
        return None
    write_cache(path, text)
    stats["extract_ok"][handler] = stats["extract_ok"].get(handler, 0) + 1
    return text


# ---------------------------------------------------------------------------
# Match helpers
# ---------------------------------------------------------------------------
def make_matcher(query: str, regex: bool):
    """Return a function that returns True if query matches a line."""
    if regex:
        try:
            rx = re.compile(query, re.IGNORECASE)
        except re.error as e:
            print(f"Invalid regex: {e}", file=sys.stderr)
            sys.exit(2)
        return lambda line: rx.search(line) is not None
    q_lower = query.lower()
    return lambda line: q_lower in line.lower()


def scan_text(text: str, matcher, context: int, max_hits: int):
    """Yield (lineno, line, context_lines) for each match."""
    lines = text.splitlines()
    hits = 0
    for i, line in enumerate(lines):
        if matcher(line):
            start = max(0, i - context)
            end = min(len(lines), i + context + 1)
            ctx = lines[start:end]
            yield (i + 1, line, ctx)
            hits += 1
            if hits >= max_hits:
                return


# ---------------------------------------------------------------------------
# Filesystem walk with exclusions
# ---------------------------------------------------------------------------
def is_excluded_repo_path(rel: str) -> bool:
    """Check if a path (relative to repo root) is excluded."""
    parts = rel.split(os.sep)
    for ex in REPO_EXCLUDE_DIRS:
        ex_parts = ex.split("/")
        # Match anywhere in path as a prefix segment
        for i in range(len(parts) - len(ex_parts) + 1):
            if parts[i:i + len(ex_parts)] == ex_parts:
                return True
    if rel in REPO_EXCLUDE_FILES:
        return True
    for suf in REPO_EXCLUDE_SUFFIXES:
        if rel.endswith(suf):
            return True
    return False


def walk_files(root: Path, exclude_fn=None):
    """Yield Path objects under root, applying exclusion."""
    for dirpath, dirnames, filenames in os.walk(root):
        # In-place prune of dirnames for efficiency
        rel_dir = os.path.relpath(dirpath, root)
        if rel_dir == ".":
            rel_dir = ""
        pruned = []
        for d in dirnames:
            sub_rel = os.path.join(rel_dir, d) if rel_dir else d
            if exclude_fn and exclude_fn(sub_rel):
                continue
            if d == ".git" or d == "node_modules":
                continue
            pruned.append(d)
        dirnames[:] = pruned
        for f in filenames:
            sub_rel = os.path.join(rel_dir, f) if rel_dir else f
            if exclude_fn and exclude_fn(sub_rel):
                continue
            yield Path(dirpath) / f


# ---------------------------------------------------------------------------
# Surface searchers
# ---------------------------------------------------------------------------
def search_filesystem(
    root: Path,
    label: str,
    matcher,
    context: int,
    limit: int,
    tools: dict,
    stats: dict,
    exclude_fn=None,
    show_relative_to: Path | None = None,
):
    """Search all files under root. Returns list of formatted hit strings."""
    hits = []
    file_count = 0
    for path in walk_files(root, exclude_fn=exclude_fn):
        if len(hits) >= limit:
            break
        ext = path.suffix.lower()
        if ext in BINARY_SUFFIXES:
            continue
        try:
            if ext in PLAIN_TEXT_EXTS or ext == "":
                # Try plain text read; skip if huge (>20MB) or undecodable
                try:
                    if path.stat().st_size > 20 * 1024 * 1024:
                        continue
                except OSError:
                    continue
                try:
                    text = path.read_text(encoding="utf-8", errors="replace")
                except (OSError, UnicodeDecodeError):
                    continue
                file_count += 1
                for lineno, line, ctx in scan_text(text, matcher, context, max_hits=5):
                    rel = path.relative_to(show_relative_to) if show_relative_to else path
                    hit = format_file_hit(str(rel), lineno, line, ctx, context)
                    hits.append(hit)
                    if len(hits) >= limit:
                        break
            elif ext in EXTRACT_EXTS:
                text = get_extracted_text(path, tools, stats)
                if text is None:
                    continue
                file_count += 1
                for lineno, line, ctx in scan_text(text, matcher, context, max_hits=3):
                    rel = path.relative_to(show_relative_to) if show_relative_to else path
                    hit = format_file_hit(
                        f"{rel} (extracted)", lineno, line, ctx, context,
                    )
                    hits.append(hit)
                    if len(hits) >= limit:
                        break
        except Exception as e:
            stats["errors"].append(f"{label}: {path}: {e}")
    stats["files_scanned"][label] = file_count
    return hits


def format_file_hit(path_str: str, lineno: int, line: str, ctx: list, context: int) -> str:
    """Format a single file hit. Trim very long lines."""
    line = line.rstrip()
    if len(line) > 240:
        line = line[:237] + "..."
    if context == 0:
        return f"  {path_str}:{lineno} — {line}"
    # Multi-line with context
    out = [f"  {path_str}:{lineno} — {line}"]
    return "\n".join(out)


def search_confluence(query: str, regex: bool, limit: int, stats: dict) -> tuple[list, str | None]:
    """Search Confluence via REST API. Returns (hits, error_or_None)."""
    if not CONFLUENCE_CREDS.exists():
        return [], "creds file missing"
    try:
        creds = json.loads(CONFLUENCE_CREDS.read_text())
    except Exception as e:
        return [], f"creds unreadable: {e}"

    site = creds.get("site", "").rstrip("/")
    email = creds.get("email") or creds.get("username")
    token = creds.get("token") or creds.get("api_token") or creds.get("password")
    if not (site and email and token):
        return [], "creds incomplete"

    auth = base64.b64encode(f"{email}:{token}".encode()).decode()
    headers = {"Authorization": f"Basic {auth}", "Accept": "application/json"}

    # CQL: text ~ "query" — Confluence handles its own search semantics.
    # For regex queries we strip the regex syntax for the CQL stage,
    # then post-filter results client-side using the regex on the body.
    if regex:
        # Use a degraded substring guess from the regex literal parts
        cql_query = re.sub(r"[\\^$.|?*+()\[\]{}]", " ", query).strip() or query
    else:
        cql_query = query

    # Escape double-quotes for CQL
    safe = cql_query.replace('"', '\\"')
    cql = f'text ~ "{safe}"'
    params = urllib.parse.urlencode({
        "cql": cql,
        "limit": str(min(limit, 25)),
        "expand": "body.view",
    })
    url = f"{site}/wiki/rest/api/content/search?{params}"

    try:
        req = urllib.request.Request(url, headers=headers)
        with urllib.request.urlopen(req, timeout=CONFLUENCE_TIMEOUT) as resp:
            data = json.loads(resp.read().decode("utf-8", errors="replace"))
    except Exception as e:
        return [], f"API call failed: {e}"

    results = data.get("results", []) or []
    matcher = make_matcher(query, regex)
    hits = []
    for item in results:
        page_id = item.get("id", "?")
        title = item.get("title", "(no title)")
        # Build URL
        links = item.get("_links", {}) or {}
        webui = links.get("webui", "")
        base = data.get("_links", {}).get("base", site + "/wiki")
        page_url = f"{base}{webui}" if webui else f"{site}/wiki/spaces/?pageId={page_id}"
        # Body excerpt: strip HTML tags lightly
        body_view = ((item.get("body") or {}).get("view") or {}).get("value", "")
        body_text = re.sub(r"<[^>]+>", " ", body_view)
        body_text = re.sub(r"\s+", " ", body_text).strip()
        # Find first matching window
        excerpt = ""
        if body_text:
            for line in body_text.split(". "):
                if matcher(line):
                    excerpt = line.strip()
                    if len(excerpt) > 200:
                        excerpt = excerpt[:197] + "..."
                    break
            if not excerpt:
                # Title-only hit — still report
                excerpt = (body_text[:160] + "...") if len(body_text) > 160 else body_text
        hits.append(
            f"  [{page_id}] {title} — {page_url}\n    {excerpt}"
        )
        if len(hits) >= limit:
            break
    stats["confluence_pages_returned"] = len(results)
    return hits, None


# ---------------------------------------------------------------------------
# Cache-warm helper (--build-cache)
# ---------------------------------------------------------------------------
def build_cache(tools: dict, stats: dict):
    """Walk Shared dir and pre-extract all docx/pdf/pptx/xlsx into cache."""
    print(f"Building extraction cache under {CACHE_DIR}...", file=sys.stderr)
    count = 0
    for path in walk_files(SHARED_DIR):
        ext = path.suffix.lower()
        if ext in EXTRACT_EXTS:
            text = get_extracted_text(path, tools, stats)
            count += 1
            if count % 50 == 0:
                print(f"  ... {count} files processed", file=sys.stderr)
    print(f"Cache build done. {count} files processed.", file=sys.stderr)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    parser = argparse.ArgumentParser(
        description="Unified search across Hexworth knowledge surfaces.",
    )
    parser.add_argument("query", nargs="?", help="search query (substring or regex)")
    parser.add_argument(
        "--surfaces",
        default="repo,memory,shared,confluence",
        help="comma-separated subset of: repo,memory,shared,confluence",
    )
    parser.add_argument("--regex", action="store_true", help="treat query as regex")
    parser.add_argument("--no-confluence", action="store_true", help="skip Confluence")
    parser.add_argument("--limit", type=int, default=50, help="max hits per surface")
    parser.add_argument("--context", type=int, default=0, help="context lines per hit")
    parser.add_argument("--build-cache", action="store_true", help="pre-extract Shared files")
    args = parser.parse_args()

    tools = check_tools()
    stats = {
        "files_scanned": {},
        "extract_ok": {},
        "extract_failed": {},
        "cache_hits": 0,
        "errors": [],
        "confluence_pages_returned": 0,
    }

    # Build-cache mode: no query needed
    if args.build_cache:
        build_cache(tools, stats)
        return 0

    if not args.query:
        parser.error("query is required (unless --build-cache)")

    surfaces = [s.strip() for s in args.surfaces.split(",") if s.strip()]
    if args.no_confluence and "confluence" in surfaces:
        surfaces.remove("confluence")

    matcher = make_matcher(args.query, args.regex)
    t0 = time.time()
    surface_errored = False
    total_hits = 0
    output_blocks = []

    # ---- Repo ----
    if "repo" in surfaces:
        hits = search_filesystem(
            REPO_ROOT, "Repo", matcher, args.context, args.limit,
            tools, stats, exclude_fn=is_excluded_repo_path,
            show_relative_to=REPO_ROOT,
        )
        output_blocks.append(("Repo", hits))
        total_hits += len(hits)

    # ---- Memory ----
    if "memory" in surfaces:
        hits = search_filesystem(
            MEMORY_DIR, "Memory", matcher, args.context, args.limit,
            tools, stats, show_relative_to=MEMORY_DIR,
        )
        output_blocks.append(("Memory", hits))
        total_hits += len(hits)

    # ---- Shared ----
    if "shared" in surfaces:
        hits = search_filesystem(
            SHARED_DIR, "Shared", matcher, args.context, args.limit,
            tools, stats,
        )
        output_blocks.append(("Shared", hits))
        total_hits += len(hits)

    # ---- Confluence ----
    if "confluence" in surfaces:
        c_hits, c_err = search_confluence(args.query, args.regex, args.limit, stats)
        if c_err:
            print(f"Confluence skipped: {c_err}", file=sys.stderr)
            surface_errored = True
            output_blocks.append(("Confluence", []))
        else:
            output_blocks.append(("Confluence", c_hits))
            total_hits += len(c_hits)

    # ---- Render output ----
    for label, hits in output_blocks:
        print(f"=== {label} ({len(hits)} hits) ===")
        if not hits:
            print("  (no hits)")
        else:
            for h in hits:
                print(h)
        print()
    sys.stdout.flush()

    elapsed = time.time() - t0
    cache_note = f" (cache: {stats['cache_hits']} hits)" if stats["cache_hits"] else ""
    print(f"Search took {elapsed:.1f}s{cache_note}. {total_hits} results total.", file=sys.stderr)

    if stats["errors"]:
        print(f"Errors: {len(stats['errors'])}", file=sys.stderr)
        for e in stats["errors"][:5]:
            print(f"  {e}", file=sys.stderr)
        surface_errored = True

    if surface_errored:
        return 2
    return 0 if total_hits > 0 else 1


if __name__ == "__main__":
    sys.exit(main())
